#!/usr/bin/env python3
"""
API Gateway PowerPoint Presentation Generator
Creates a comprehensive presentation explaining API Gateway concepts,
architecture, integrations, and use cases for the ECCS platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle=""):
    """Create a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    if subtitle:
        subtitle_shape.text = subtitle
    
    return slide

def create_content_slide(prs, title, content_points=None):
    """Create a slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    if content_points:
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, point in enumerate(content_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Handle nested bullets (tuples)
            if isinstance(point, tuple):
                p.text = point[0]
                p.level = 0
                for sub_point in point[1]:
                    p = text_frame.add_paragraph()
                    p.text = sub_point
                    p.level = 1
            else:
                p.text = point
                p.level = 0
    
    return slide

def create_section_header(prs, section_title):
    """Create a section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = section_title
    
    return slide

def add_code_block(slide, code_text, left, top, width, height):
    """Add a code block to a slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = code_text
    p.font.name = 'Courier New'
    p.font.size = Pt(9)
    
    # Set background color for code block
    fill = textbox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    return textbox

def main():
    """Generate the comprehensive API Gateway presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ============================================================================
    # TITLE SLIDE
    # ============================================================================
    create_title_slide(
        prs,
        "API Gateway Deep Dive",
        "Understanding Traefik in the ECCS Platform\nArchitecture, Integrations, and Best Practices"
    )
    
    # ============================================================================
    # SECTION 1: INTRODUCTION
    # ============================================================================
    create_section_header(prs, "Section 1: Introduction to API Gateways")
    
    create_content_slide(prs, "What is an API Gateway?", [
        "A single entry point for all client requests to backend services",
        "Acts as a reverse proxy that routes requests to appropriate microservices",
        "Provides cross-cutting concerns: authentication, rate limiting, monitoring",
        "Decouples clients from backend service implementations",
        "Enables microservices architecture by managing service complexity"
    ])
    
    create_content_slide(prs, "Why API Gateways?", [
        ("Problems They Solve:", [
            "Service discovery - How do clients find services?",
            "Protocol translation - HTTP, gRPC, WebSocket handling",
            "Security - Centralized authentication and authorization",
            "Performance - Load balancing, caching, compression",
            "Observability - Centralized logging and metrics"
        ]),
        ("Benefits:", [
            "Simplified client code - Single endpoint instead of many",
            "Reduced coupling - Services can change without affecting clients",
            "Better security - Single point to enforce policies",
            "Improved performance - Optimized routing and caching"
        ])
    ])
    
    create_content_slide(prs, "API Gateway in ECCS Architecture", [
        "Traefik serves as the API Gateway for ECCS platform",
        "All external traffic flows through Traefik first",
        "Routes requests to: Auth Service, Email Service, Locations Service",
        "Enforces security policies before reaching backend services",
        "Provides observability through Prometheus, Jaeger, and access logs",
        "Supports both HTTP and HTTPS with TLS termination"
    ])
    
    # ============================================================================
    # SECTION 2: TRAEFIK ARCHITECTURE
    # ============================================================================
    create_section_header(prs, "Section 2: Traefik Architecture Deep Dive")
    
    create_content_slide(prs, "Traefik Core Components", [
        ("EntryPoints:", [
            "Network entry points (ports) that listen for incoming traffic",
            "Example: web (8800), websecure (8443), traefik (8080)",
            "Each entrypoint can have specific configurations"
        ]),
        ("Routers:", [
            "Define routing rules to match incoming requests",
            "Connect entrypoints to services via middlewares",
            "Priority-based matching (lower number = higher priority)"
        ]),
        ("Services:", [
            "Define backend servers and load balancing strategy",
            "Include health checks for high availability",
            "Support multiple servers per service"
        ]),
        ("Middlewares:", [
            "Process requests/responses between router and service",
            "Examples: authentication, rate limiting, headers, compression"
        ])
    ])
    
    create_content_slide(prs, "Request Flow Through Traefik", [
        "1. Client sends request to EntryPoint (e.g., port 8800)",
        "2. Traefik matches request against Router rules",
        "3. Selected router's middlewares process the request in order:",
        "   → Rate limiting (block DDoS attacks first)",
        "   → Security headers (add protective headers)",
        "   → JWT authentication (verify user identity)",
        "   → Retry logic (handle transient failures)",
        "4. Request forwarded to backend Service",
        "5. Service performs health check if configured",
        "6. Load balancer selects healthy backend server",
        "7. Response flows back through middleware chain",
        "8. Client receives final response"
    ])
    
    create_content_slide(prs, "Dynamic Service Discovery", [
        ("File Provider:", [
            "Static configuration loaded from YAML files",
            "Changes detected automatically with hot-reload",
            "Used for core services and middleware definitions"
        ]),
        ("Docker/Podman Provider (Optional):", [
            "Automatic discovery of containerized services",
            "Services register via container labels",
            "Zero-downtime deployments with auto-discovery",
            "Requires container socket access"
        ]),
        ("ECCS Approach:", [
            "Uses File Provider for reliability",
            "All services defined in dynamic-config.yml",
            "Avoids permission issues with Podman socket"
        ])
    ])
    
    # Save presentation (continuing in next part due to character limits)
    print("Creating presentation part 1...")
    
